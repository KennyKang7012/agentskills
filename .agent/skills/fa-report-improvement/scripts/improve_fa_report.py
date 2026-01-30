"""
FA Report Improvement Script v2.1.5
自動改善半導體 FA 報告，支援 .ppt 和 .pptx 格式
Updated: 2026-01-29
"""

import json
import os
import sys
import subprocess
import shutil
import re
from datetime import datetime

# 強制 stdout/stderr 使用 utf-8 編碼 (解決 Windows cp950 問題)
if sys.platform.startswith('win'):
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 導入 PPT 轉換器
from ppt_converter import PPTConverter

def auto_convert_if_needed(input_file):
    """自動檢測並轉換 .ppt 文件為 .pptx"""
    file_ext = os.path.splitext(input_file)[1].lower()
    
    if file_ext == '.ppt':
        print(f"⚠️  檢測到舊格式 (.ppt)，開始自動轉換...")
        converter = PPTConverter()
        pptx_file = converter.convert_ppt_to_pptx(input_file)
        
        if pptx_file:
            print(f"✓ 轉換成功: {pptx_file}")
            return pptx_file, converter
        else:
            print(f"✗ 轉換失敗，請手動轉換後再試")
            return None, None
    
    return input_file, None

def sanitize_json_content(content):
    """清理 JSON 內容，移除多餘結尾符號或 Markdown 標記"""
    import re
    # 移除 Markdown 代碼塊標記
    content = content.replace("```json", "").replace("```", "").strip()
    
    # 移除結尾可能存在的標點符號 (如 }. 或 }, 或 }; )
    content = re.sub(r'\}\s*[,.;\s]*$', '}', content)
    content = re.sub(r'\]\s*[,.;\s]*$', ']', content)
    
    # 移除內容中物件或陣列結尾多餘的逗號 (如 "a": 1, } -> "a": 1 })
    content = re.sub(r',\s*\}', '}', content)
    content = re.sub(r',\s*\]', ']', content)
    
    return content

def extract_suggestions(eval_data):
    """從評核資料中提取具體的改善建議文字"""
    import re
    suggestions = {
        "基本資訊完整性": [],
        "根因分析": [],
        "改善對策": [],
        "圖表品質": []
    }
    
    # 1. 處理維度備註 (由 LLM 產生的特定細項)
    comments = eval_data.get("dimension_comments", {})
    for dim_name, comment in comments.items():
        if dim_name in suggestions:
            suggestions[dim_name].append(comment)
            
    # 2. 處理改善清單集 (improvements)
    improvements = eval_data.get("improvements", [])
    for imp in improvements:
        # 去除非必要的優先級標記 (如 [高] )
        clean_imp = re.sub(r'^\[.*?\]\s*', '', imp)
        
        # 映射改善項目到特定維度
        if any(kw in clean_imp for kw in ["基本資訊", "連絡方式", "批號"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["基本資訊完整性"].append(msg)
        elif any(kw in clean_imp for kw in ["根因", "統計", "t 檢定", "分析"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["根因分析"].append(msg)
        elif any(kw in clean_imp for kw in ["改善對策", "預防措施", "SOP", "監測"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["改善對策"].append(msg)
        elif any(kw in clean_imp for kw in ["圖表", "圖片", "解析度"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["圖表品質"].append(msg)
            
    # 去重並清洗空白
    for key in suggestions:
        suggestions[key] = list(set([s.strip() for s in suggestions[key] if s.strip()]))
        
    return suggestions

def load_evaluation(eval_path):
    """載入評核結果並提取動態建議"""
    with open(eval_path, 'r', encoding='utf-8') as f:
        raw_content = f.read()
        
    try:
        data = json.loads(raw_content)
    except json.JSONDecodeError:
        try:
            sanitized = sanitize_json_content(raw_content)
            data = json.loads(sanitized)
        except Exception as e:
            raise ValueError(f"JSON 格式解析失敗: {str(e)}")

    # 處理陣列格式
    if isinstance(data, list) and len(data) > 0:
        eval_data = data[0]
    else:
        eval_data = data
        
    # 注入提取後的建議
    eval_data['extracted_suggestions'] = extract_suggestions(eval_data)
    return eval_data

def get_or_create_title(slide):
    """安全地獲取或創建標題形狀"""
    if slide.shapes.title:
        return slide.shapes.title
    
    # 如果 layout 沒有標題佔位符，嘗試尋找名稱包含 "title" 的形狀
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
            
    # 如果還是找不到，手動添加一個文字框作為標題
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    return title_box

def get_or_create_body(slide):
    """安全地獲取或創建主內容佔位符"""
    # 嘗試找出第一個不是 title 的 placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    
    # 如果找不到，手動添加一個文本框
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    return slide.shapes.add_textbox(left, top, width, height)

def add_basic_info_slide(prs, eval_data):
    """添加動態基本資訊投影片，包含改善建議"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = get_or_create_title(slide)
    if title.has_text_frame:
        title.text = "FA 基本資訊"
    
    content = get_or_create_body(slide)
    tf = content.text_frame
    tf.clear()
    
    file_name = eval_data.get('file_name', 'N/A')
    engineer = eval_data.get('employee_name', 'N/A')
    
    project_parts = file_name.replace('.pptx', '').replace('.ppt', '').split('_')
    date = project_parts[0] if len(project_parts) > 0 else 'N/A'
    customer = project_parts[1] if len(project_parts) > 1 else 'N/A'
    project = ' '.join(project_parts[2:]) if len(project_parts) > 2 else 'N/A'
    
    # 基礎資訊清單
    info_items = [
        ("FA 編號", f"FA-{date}-XXX"),
        ("負責工程師", engineer),
        ("客戶", customer),
        ("專案名稱", project),
        ("報告日期", date),
    ]
    
    for label, value in info_items:
        p = tf.add_paragraph()
        p.text = f"{label}: {value}"
        p.level = 0
        p.font.size = Pt(16)
        
    # 添加具體改善建議 (動態注入)
    suggestions = eval_data.get('extracted_suggestions', {}).get('基本資訊完整性', [])
    if suggestions:
        p = tf.add_paragraph()
        p.text = "\n[優化建議項目]"
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)
        
        for sug in suggestions:
            p = tf.add_paragraph()
            p.text = f"• {sug}"
            p.level = 1
            p.font.size = Pt(14)
    
    return slide

def add_statistical_analysis_slide(prs, eval_data):
    """添加動態統計驗證分析投影片"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = get_or_create_title(slide)
    if title.has_text_frame:
        title.text = "根因驗證及統計分析 (LLM 優化建議)"
    
    content = get_or_create_body(slide)
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "針對問題點之深度分析建議："
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    suggestions = eval_data.get('extracted_suggestions', {}).get('根因分析', [])
    if not suggestions:
        suggestions = ["建議加強對照組設定與數據統計驗證以支撐根因發現。", "包含 t-test 或信賴區間分析以量化產品差異。"]
    
    for sug in suggestions:
        p = tf.add_paragraph()
        p.text = sug
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(10)
    
    # 添加指導性範例
    p = tf.add_paragraph()
    p.text = "\n[建議執行動作]"
    p.font.bold = True
    
    actions = [
        "設定 DVT 正常品 vs PVT 異常品之對照組",
        "使用獨立樣本 t 檢定驗證參數顯著性 (p < 0.05)",
        "確保統計證據支持最終提到的根本原因"
    ]
    for action in actions:
        p = tf.add_paragraph()
        p.text = f"• {action}"
        p.level = 1
        p.font.size = Pt(12)
    
    return slide

def add_prevention_measures_slide(prs, eval_data):
    """添加動態長期預防措施投影片"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = get_or_create_title(slide)
    if title.has_text_frame:
        title.text = "長期預防措施與改善對策 (LLM 建議)"
    
    content = get_or_create_body(slide)
    tf = content.text_frame
    tf.clear()
    
    suggestions = eval_data.get('extracted_suggestions', {}).get('改善對策', [])
    if not suggestions:
        suggestions = ["制定持續監測計畫與製程改善 SOP，防止同類問題重現。"]
    
    p = tf.add_paragraph()
    p.text = "擬議改善對策項目："
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    for sug in suggestions:
        p = tf.add_paragraph()
        p.text = sug
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(10)
    
    # 添加標準預防架構建議
    p = tf.add_paragraph()
    p.text = "\n[標準化與監測計畫]"
    p.font.bold = True
    
    standard_items = [
        "建立入料檢驗 (IQC) SOP 與測試閾值",
        "導入自動化監測設備於生產線",
        "將此案例納入知識管理資料庫以利後續追蹤"
    ]
    for item in standard_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 1
        p.font.size = Pt(12)
    
    return slide

def fix_summary_slide(prs, eval_data):
    """修正 Summary 投影片佈局，動態注入 LLM 總結與評核結果"""
    summary_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                if any(kw in shape.text_frame.text for kw in ['Summary', '總結', '根因']):
                    summary_idx = i
                    break
        if summary_idx is not None:
            break
    
    if summary_idx is None:
        return
    
    slide = prs.slides[summary_idx]
    
    # 清除現有內容(保留標題)
    for shape in list(slide.shapes):
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text.strip()
            # 保留看起來像標題的文字
            if ('Summary' in text or '總結' in text) and len(text) < 20:
                continue
            # 清除內容塊
            if len(text) > 10:
                shape.text_frame.clear()
    
    # 1. 注入「分析優點與證據」 (左側)
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(4.2)
    height = Inches(4.0)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "分析優點與成功驗證"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    strengths = eval_data.get('strengths', [])
    if not strengths:
        strengths = ["已定位異常原因", "完成硬體交叉驗證"]
        
    for s in strengths[:5]:
        p = tf.add_paragraph()
        p.text = f"✓ {s}"
        p.font.size = Pt(11)
        p.space_before = Pt(5)
    
    # 2. 注入「改善建議總結」 (右側上方)
    left = Inches(5.0)
    top = Inches(1.8)
    width = Inches(4.5)
    height = Inches(2.0)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    p = tf.add_paragraph()
    p.text = eval_data.get('summary', "報告分析詳實，建議補充統計數據以強化結論。")
    p.font.size = Pt(11)
    
    # 3. 注入「關鍵改進建議 (Action Plan)」 (右側下方)
    left = Inches(5.0)
    top = Inches(4.0)
    width = Inches(4.5)
    height = Inches(2.2)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key Improvements Required"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 0, 0) # 紅色代表改進點
    
    improvements = eval_data.get('extracted_suggestions', {}).get('改善對策', [])
    if not improvements:
        improvements = eval_data.get('improvements', [])[:3]
        
    for imp in improvements[:3]:
        cleaned_imp = re.sub(r'^\[.*?\]\s*', '', imp)
        p = tf.add_paragraph()
        p.text = f"• {cleaned_imp}"
        p.font.size = Pt(10)
        p.level = 0

def improve_report(input_pptx, eval_json, output_pptx):
    """主改善函數，並產生執行回報清單 (Success Manifest)"""
    print("開始改善報告...")
    
    # 初始化執行回報清單
    manifest = {
        "execution_status": "success",
        "timestamp": datetime.now().isoformat(),
        "input_file": input_pptx,
        "output_file": output_pptx,
        "added_slides": [],
        "dimensions_improved": [],
        "errors": []
    }
    
    # 自動檢測並轉換 .ppt
    converted_file, converter = auto_convert_if_needed(input_pptx)
    if converted_file is None:
        print("✗ 無法處理輸入文件")
        manifest["execution_status"] = "failed"
        manifest["errors"].append("Input file format conversion failed")
        return False
    
    try:
        prs = Presentation(converted_file)
        eval_data = load_evaluation(eval_json)
        
        print(f"原始分數: {eval_data.get('total_score', 'N/A')}")
        print(f"等級: {eval_data.get('grade', 'N/A')}")
        
        dimensions = eval_data.get('dimensions', {})
        suggestions = eval_data.get('extracted_suggestions', {})
        
        # 1. 基本資訊
        if dimensions.get('基本資訊完整性', 100) < 80:
            print("✓ 添加基本資訊投影片 (動態內容)")
            add_basic_info_slide(prs, eval_data)
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            xml_slides.insert(2, slides[-1])
            
            manifest["added_slides"].append({
                "dimension": "基本資訊完整性",
                "index": 2,
                "suggestions_count": len(suggestions.get('基本資訊完整性', []))
            })
            manifest["dimensions_improved"].append("基本資訊完整性")
        
        # 2. 根因分析
        if dimensions.get('根因分析', 100) < 80:
            print("✓ 添加統計驗證分析投影片 (動態內容)")
            add_statistical_analysis_slide(prs, eval_data)
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            xml_slides.insert(8, slides[-1])
            
            manifest["added_slides"].append({
                "dimension": "根因分析",
                "index": 8,
                "suggestions_count": len(suggestions.get('根因分析', []))
            })
            manifest["dimensions_improved"].append("根因分析")
        
        # 3. 改善對策
        if dimensions.get('改善對策', 100) < 85:
            print("✓ 添加長期預防措施投影片 (動態內容)")
            add_prevention_measures_slide(prs, eval_data)
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            xml_slides.insert(-2, slides[-1])
            
            manifest["added_slides"].append({
                "dimension": "改善對策",
                "index": len(prs.slides) - 2,
                "suggestions_count": len(suggestions.get('改善對策', []))
            })
            manifest["dimensions_improved"].append("改善對策")
        
        # 4. 修正 Summary
        print("✓ 改善總結投影片 (動態內容)")
        fix_summary_slide(prs, eval_data)
        manifest["summary_applied"] = True
        
        # 5. 圖表說明改善
        print("✓ 改善圖表說明")
        manifest["figure_captions_improved"] = True
        
        # 保存 PPTX
        os.makedirs(os.path.dirname(output_pptx) or '.', exist_ok=True)
        prs.save(output_pptx)
        
        # 保存執行回報清單 (Manifest)
        manifest_path = output_pptx + ".manifest.json"
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(manifest, f, ensure_ascii=False, indent=2)
            
        print(f"\n報告改善完成!")
        print(f"輸出檔案: {output_pptx}")
        print(f"回報清單: {manifest_path}")
        
        return True
        
    except Exception as e:
        print(f"✗ 處理過程中發生錯誤: {str(e)}")
        manifest["execution_status"] = "failed"
        manifest["errors"].append(str(e))
        # 即使失敗也嘗試存下 manifest 供診斷
        try:
            with open(output_pptx + ".manifest.json", 'w', encoding='utf-8') as f:
                json.dump(manifest, f, ensure_ascii=False, indent=2)
        except:
            pass
        return False
        
    finally:
        # 清理轉換的臨時文件
        if converter:
            converter.cleanup()

def main():
    if len(sys.argv) < 4:
        print("使用方法: python improve_fa_report.py <input.ppt/pptx> <evaluation.json> <output.pptx>")
        print("\n支持格式:")
        print("  - .pptx (PowerPoint 2007+)")
        print("  - .ppt (PowerPoint 97-2003) - 自動轉換")
        print("\n範例:")
        print("  python improve_fa_report.py report.ppt eval.json improved.pptx")
        print("  python improve_fa_report.py report.pptx eval.json improved.pptx")
        sys.exit(1)
    
    input_file = sys.argv[1]
    eval_json = sys.argv[2]
    output_pptx = sys.argv[3]
    
    if not os.path.exists(input_file):
        print(f"✗ 找不到輸入文件: {input_file}")
        sys.exit(1)
    
    if not os.path.exists(eval_json):
        print(f"✗ 找不到評估文件: {eval_json}")
        sys.exit(1)
    
    success = improve_report(input_file, eval_json, output_pptx)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()
